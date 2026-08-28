#!/usr/bin/env python3
"""15 dakikalık bitirme sunumu (PowerPoint). Streamlit'e dokunmaz, modeli yeniden eğitmez."""

from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG = ROOT / "reports" / "bitirme_figures"
SHAP_CSV = ROOT / "reports" / "explainability" / "shap_feature_importance.csv"
ASSET = ROOT / "reports" / "pptx_figures"
OUT = ROOT / "docs" / "Elektrik_Piyasasi_Fiyat_Tahmini_Sunum.pptx"

NAVY = RGBColor(0x1F, 0x4E, 0x9B)
NAVY_DARK = RGBColor(0x1F, 0x2A, 0x44)
MUTED = RGBColor(0x5B, 0x6B, 0x84)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PALE = RGBColor(0xF3, 0xF6, 0xFB)
LINE = RGBColor(0xD7, 0xDE, 0xEA)
GREEN = RGBColor(0x16, 0x65, 0x34)
GREEN_BG = RGBColor(0xEA, 0xF6, 0xEE)
WARN = RGBColor(0x9B, 0x1C, 0x1C)
WARN_BG = RGBColor(0xFD, 0xEC, 0xEC)

WF = [
    ("Naive Lag-24", 7.351456),
    ("Ridge", 5.796612),
    ("LightGBM", 5.916945),
    ("XGBoost", 5.945953),
    ("Ridge+B+AR(1)", 4.529617),
]

FRIENDLY = {
    "price_day_ahead_lag_24": "Fiyat gecikmesi 24s",
    "price_day_ahead_lag_48": "Fiyat gecikmesi 48s",
    "price_day_ahead_lag_168": "Fiyat gecikmesi 168s",
    "price_mean_lag24_lag48": "Fiyat ort. gecikme 24–48s",
    "total_load_forecast": "Toplam yük tahmini",
    "forecast_solar_day_ahead": "Güneş tahmini",
    "forecast_wind_onshore_day_ahead": "Rüzgâr tahmini",
    "renewable_forecast_total": "Yenilenebilir tahmin toplamı",
    "forecast_wind_share_of_load": "Rüzgârın yük payı",
    "generation_wind_onshore_lag_24": "Rüzgâr üretimi gecikme 24s",
}


def set_run(run, *, size=18, bold=False, color=NAVY_DARK, font="Calibri"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def add_text(tf, text, *, size=18, bold=False, color=NAVY_DARK, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def add_para(tf, text, *, size=16, bold=False, color=NAVY_DARK, space_before=6, space_after=4):
    p = tf.add_paragraph()
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, color=color)
    return p


def fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill_shape(sh, color)
    sh.adjustments[0] = 0.08
    return sh


def notes(slide, text: str) -> None:
    slide.notes_slide.notes_text_frame.text = text


def footer(slide, n: int, total: int = 10) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.15), Inches(11.5), Inches(0.28))
    add_text(box.text_frame, f"Elektrik Piyasası Fiyat Tahmini   ·   {n} / {total}", size=11, color=MUTED)


def kicker(slide, text: str) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.12))
    fill_shape(bar, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.28), Inches(12), Inches(0.32))
    add_text(box.text_frame, text.upper(), size=12, bold=True, color=MUTED)


def title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.52), Inches(12.1), Inches(0.55))
    add_text(box.text_frame, text, size=28, bold=True, color=NAVY_DARK)


def card(slide, l, t, w, h, label: str, value: str, *, hl=False):
    sh = rect(slide, l, t, w, h, PALE if hl else WHITE)
    sh.line.fill.solid()
    sh.line.color.rgb = NAVY if hl else LINE
    sh.line.width = Pt(1.25)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = label
    set_run(r, size=12, color=MUTED)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(4)
    r2 = p2.add_run()
    r2.text = value
    set_run(r2, size=22, bold=True, color=NAVY if hl else NAVY_DARK)


def make_wf_chart(path: Path) -> None:
    names = [x[0] for x in WF]
    vals = [x[1] for x in WF]
    colors = ["#8AA0C2"] * 4 + ["#1F4E9B"]
    fig, ax = plt.subplots(figsize=(10.8, 4.2), dpi=160)
    bars = ax.bar(names, vals, color=colors, width=0.62)
    ax.set_ylabel("Walk-forward MAE (€/MWh)", fontsize=11, color="#1F2A44")
    ax.set_ylim(0, 8.4)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(colors="#1F2A44")
    ax.yaxis.grid(True, linestyle="--", alpha=0.35)
    ax.set_axisbelow(True)
    for b, v in zip(bars, vals):
        ax.text(b.get_x() + b.get_width() / 2, v + 0.18, f"{v:.2f}", ha="center", va="bottom", fontsize=11, fontweight="600", color="#1F2A44")
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def make_shap_chart(path: Path) -> None:
    wanted = [
        "price_day_ahead_lag_24",
        "price_day_ahead_lag_48",
        "total_load_forecast",
        "forecast_wind_share_of_load",
        "forecast_solar_day_ahead",
        "renewable_forecast_total",
        "generation_wind_onshore_lag_24",
        "price_mean_lag24_lag48",
        "price_day_ahead_lag_168",
        "forecast_wind_onshore_day_ahead",
    ]
    df = pd.read_csv(SHAP_CSV)
    show = df[df["feature"].isin(wanted)].copy()
    show["label"] = show["feature"].map(FRIENDLY).fillna(show["feature"])
    show = show.sort_values("mean_abs_shap")
    fig, ax = plt.subplots(figsize=(10.2, 4.6), dpi=160)
    ax.barh(show["label"], show["mean_abs_shap"], color="#1F4E9B")
    ax.set_xlabel("Ort. |SHAP|", fontsize=11, color="#1F2A44")
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    ax.tick_params(colors="#1F2A44")
    ax.xaxis.grid(True, linestyle="--", alpha=0.35)
    ax.set_axisbelow(True)
    fig.tight_layout()
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def build() -> Path:
    ASSET.mkdir(parents=True, exist_ok=True)
    wf_png = ASSET / "wf_mae.png"
    shap_png = ASSET / "shap_top.png"
    make_wf_chart(wf_png)
    make_shap_chart(shap_png)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    # 1 Title
    s = prs.slides.add_slide(blank)
    top = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    fill_shape(top, NAVY_DARK)
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), Inches(7.5))
    fill_shape(accent, NAVY)
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.05), Inches(11.4), Inches(0.4))
    add_text(box.text_frame, "BİTİRME PROJESİ SUNUMU", size=14, bold=True, color=RGBColor(0x8A, 0xA0, 0xC2))
    box = s.shapes.add_textbox(Inches(0.9), Inches(2.45), Inches(11.5), Inches(1.3))
    add_text(box.text_frame, "Elektrik Piyasası Fiyat Tahmini", size=40, bold=True, color=WHITE)
    box = s.shapes.add_textbox(Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.55))
    add_text(box.text_frame, "İspanya gün öncesi piyasası  ·  saatlik fiyat (€/MWh)  ·  2014–2018", size=20, color=RGBColor(0xC5, 0xD0, 0xE0))
    box = s.shapes.add_textbox(Inches(0.9), Inches(5.35), Inches(11.5), Inches(0.4))
    add_text(box.text_frame, "Ridge (α = 0.001)  +  METHOD_B  +  AR(1)", size=18, bold=True, color=WHITE)
    logo = FIG / "logo_softito.png"
    if logo.exists():
        s.shapes.add_picture(str(logo), Inches(10.55), Inches(0.35), height=Inches(0.72))
    notes(
        s,
        "Bu çalışmada İspanya elektrik piyasasının saatlik gün öncesi fiyatı tahmin ediliyor. "
        "Karşılaştırma olarak dünün aynı saatindeki fiyat kullanılıyor.",
    )

    # 2 Problem
    s = prs.slides.add_slide(blank)
    kicker(s, "01  ·  Problem ve amaç")
    title(s, "Ne tahmin ediliyor?")
    q = rect(s, Inches(0.6), Inches(1.2), Inches(12.1), Inches(1.15), PALE)
    q.line.fill.solid()
    q.line.color.rgb = NAVY
    tf = q.text_frame
    tf.word_wrap = True
    add_text(tf, "Model, sadece dünün aynı saatindeki fiyatı kullanmaktan daha iyi tahmin yapabilir mi?", size=22, bold=True, color=NAVY)
    card(s, Inches(0.6), Inches(2.6), Inches(3.9), Inches(1.25), "Hedef", "Gün öncesi fiyat (€/MWh)", hl=True)
    card(s, Inches(4.7), Inches(2.6), Inches(2.55), Inches(1.25), "Gözlem", "35.064 saat")
    card(s, Inches(7.45), Inches(2.6), Inches(2.55), Inches(1.25), "Dönem", "2014–2018")
    card(s, Inches(10.2), Inches(2.6), Inches(2.5), Inches(1.25), "Piyasa", "İspanya")
    box = s.shapes.add_textbox(Inches(0.6), Inches(4.15), Inches(12.1), Inches(2.5))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Bu çalışmada İspanya elektrik piyasasının saatlik gün öncesi fiyatı tahmin ediliyor. "
        "Karşılaştırma olarak dünün aynı saatindeki fiyat kullanılıyor. "
        "Geçmiş fiyat, tüketim, üretim, hava ve gün öncesi tahminlerle sızıntısız bir zaman serisi modeli kuruluyor "
        "ve bu basit kopyalamadan daha iyi sonuç alınıp alınmadığına bakılıyor.",
        size=18,
        color=NAVY_DARK,
    )
    footer(s, 2)
    notes(s, box.text_frame.paragraphs[0].runs[0].text)

    # 3 Data
    s = prs.slides.add_slide(blank)
    kicker(s, "02  ·  Veri ve özellikler")
    title(s, "Hangi veriler kullanıldı?")
    card(s, Inches(0.6), Inches(1.2), Inches(3.9), Inches(1.1), "Enerji gözlemi", "35.064")
    card(s, Inches(4.7), Inches(1.2), Inches(3.9), Inches(1.1), "Hava kaydı", "178.396")
    card(s, Inches(8.8), Inches(1.2), Inches(3.9), Inches(1.1), "Sızıntısız özellik", "184+")
    labels = ["Üretim", "Yük / Tüketim", "Hava", "Gün öncesi tahmin", "Piyasa fiyatı"]
    for i, lab in enumerate(labels):
        c = rect(s, Inches(0.6 + i * 2.46), Inches(2.5), Inches(2.32), Inches(0.7), WHITE)
        c.line.fill.solid()
        c.line.color.rgb = LINE
        add_text(c.text_frame, lab, size=14, bold=True, color=NAVY_DARK, align=PP_ALIGN.CENTER)
        c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    box = s.shapes.add_textbox(Inches(0.6), Inches(3.4), Inches(12.1), Inches(3.3))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Enerji piyasası ve hava verisi birleştirildi. Yalnızca gün öncesi ihalenin yapıldığı anda elde olan bilgiler kullanıldı. "
        "Aynı saatin gerçekleşen tüketimi ve gerçekleşen fiyatı özellik olarak alınmadı.",
        size=17,
        color=NAVY_DARK,
    )
    add_para(
        tf,
        "Üretim ve yük, arz ile talebi taşıdığı için; hava, rüzgâr ve güneşi etkilediği için; resmi gün öncesi tahminler de operatörün elindeki bilgi olduğu için kullanıldı.",
        size=17,
        space_before=12,
    )
    add_para(
        tf,
        "Fiyat gecikmesi olarak dün, iki gün önce ve geçen haftanın aynı saati alındı (t-24 · t-48 · t-168). Bir saat önceki fiyat kullanılmadı, çünkü ihale anında o fiyat henüz yoktu.",
        size=17,
        space_before=8,
    )
    footer(s, 3)
    notes(
        s,
        "Enerji piyasası ve hava verisi birleştirildi. Yalnızca gün öncesi ihalenin yapıldığı anda elde olan bilgiler kullanıldı. "
        "Aynı saatin gerçekleşen tüketimi ve gerçekleşen fiyatı özellik olarak alınmadı. "
        "Üretim ve yük, arz ile talebi taşıdığı için; hava, rüzgâr ve güneşi etkilediği için; "
        "resmi gün öncesi tahminler de operatörün elindeki bilgi olduğu için kullanıldı. "
        "Fiyat gecikmesi olarak dün, iki gün önce ve geçen haftanın aynı saati alındı. "
        "Bir saat önceki fiyat kullanılmadı, çünkü ihale anında o fiyat henüz yoktu.",
    )

    # 4 Models chart
    s = prs.slides.add_slide(blank)
    kicker(s, "03  ·  Model seçimi")
    title(s, "Hangi model seçildi?")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Modeller, test seti açılmadan walk-forward ortalama mutlak hata ile karşılaştırıldı. Düşük çubuk daha iyidir. "
        "Naive tabanın yanında Ridge, LightGBM ve XGBoost denendi. En düşük hata Ridge tabanlı birleşik modelde alındığı için bu model seçildi.",
        size=16,
        color=NAVY_DARK,
    )
    s.shapes.add_picture(str(wf_png), Inches(0.7), Inches(2.05), width=Inches(11.9))
    footer(s, 4)
    notes(
        s,
        "Modeller, test seti açılmadan walk-forward ortalama mutlak hata ile karşılaştırıldı. "
        "Düşük çubuk daha iyidir. Naive tabanın yanında Ridge, LightGBM ve XGBoost denendi. "
        "En düşük hata Ridge tabanlı birleşik modelde alındığı için bu model seçildi.",
    )

    # 5 Three components
    s = prs.slides.add_slide(blank)
    kicker(s, "03  ·  Model seçimi")
    title(s, "Seçilen model: Ridge + METHOD_B + AR(1)")
    defs = [
        (
            "Ridge (α = 0.001)",
            "Ana regresyon modeli olarak kullanılıyor. Çoklu doğrusal bağlantının etkisini azaltarak üretim, tüketim ve hava durumu gibi değişkenlerden fiyat tahmini yapıyor.",
        ),
        (
            "METHOD_B",
            "Ridge’in yakalayamadığı ek yapı modellenecek şekilde ekleniyor; tahmin bu bileşenle iyileştiriliyor.",
        ),
        (
            "AR(1)",
            "Ridge artığındaki bir önceki zaman adımı kullanılarak zamansal bağımlılık yakalanıyor.",
        ),
    ]
    for i, (h, body) in enumerate(defs):
        sh = rect(s, Inches(0.6), Inches(1.25 + i * 1.75), Inches(12.1), Inches(1.58), WHITE)
        sh.line.fill.solid()
        sh.line.color.rgb = NAVY if i == 0 else LINE
        left = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.25 + i * 1.75), Inches(0.12), Inches(1.58))
        fill_shape(left, NAVY)
        tb = s.shapes.add_textbox(Inches(1.0), Inches(1.35 + i * 1.75), Inches(11.4), Inches(1.35))
        tf = tb.text_frame
        tf.word_wrap = True
        add_text(tf, h, size=20, bold=True, color=NAVY)
        add_para(tf, body, size=17, space_before=8)
    footer(s, 5)
    notes(
        s,
        "Ridge ana regresyon modeli olarak kullanılıyor. Çoklu doğrusal bağlantının etkisini azaltarak üretim, tüketim ve hava durumu gibi değişkenlerden fiyat tahmini yapıyor. "
        "METHOD_B ile Ridge’in yakalayamadığı ek yapı modellenecek şekilde ekleniyor; tahmin bu bileşenle iyileştiriliyor. "
        "AR(1) ile Ridge artığındaki bir önceki zaman adımı kullanılarak zamansal bağımlılık yakalanıyor.",
    )

    # 6 Locked test
    s = prs.slides.add_slide(blank)
    kicker(s, "04  ·  Kilitli test")
    title(s, "Kilitli sonuç")
    card(s, Inches(0.6), Inches(1.15), Inches(3.9), Inches(1.15), "MAE", "3,99 €/MWh", hl=True)
    card(s, Inches(4.7), Inches(1.15), Inches(3.9), Inches(1.15), "Naive Lag-24", "6,05 €/MWh")
    card(s, Inches(8.8), Inches(1.15), Inches(3.9), Inches(1.15), "MAE iyileşmesi", "≈ %34")
    box = s.shapes.add_textbox(Inches(0.6), Inches(2.4), Inches(12.1), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Test seti yalnızca bir kez açıldı. Seçilen modelin ortalama mutlak hatası 3,99 euro, Naive tabanınki 6,05 euro; hata yaklaşık yüzde 34 daha düşük. "
        "Bu grafikte gerçek fiyatlarla tahmin karşılaştırılıyor. Henüz oluşmamış gelecek saatler için tahmin üretilmedi.",
        size=15,
        color=NAVY_DARK,
    )
    pred_png = FIG / "fig_actual_pred.png"
    if pred_png.exists():
        s.shapes.add_picture(str(pred_png), Inches(0.7), Inches(3.25), width=Inches(11.9))
    footer(s, 6)
    notes(
        s,
        "Test seti yalnızca bir kez açıldı. Seçilen modelin ortalama mutlak hatası 3,99 euro, Naive tabanınki 6,05 euro; hata yaklaşık yüzde 34 daha düşük. "
        "Bu grafikte gerçek fiyatlarla tahmin karşılaştırılıyor. Henüz oluşmamış gelecek saatler için tahmin üretilmedi.",
    )

    # 7 Errors
    s = prs.slides.add_slide(blank)
    kicker(s, "05  ·  Hata analizi")
    title(s, "Model nerede hata yapıyor?")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Artık, tahminden gerçeğin çıkarılmasıyla hesaplanıyor. Pozitif değer aşırı tahmini, negatif değer eksik tahmini gösteriyor. "
        "Hata rastgele dağılmıyor. Test döneminde model ortalama olarak biraz yüksek tahmin ediyor ve yüksek fiyatlı saatlerde sapma devam ediyor. "
        "Bu yüzden model üretime hazır sayılmıyor.",
        size=15,
        color=NAVY_DARK,
    )
    r_png = FIG / "fig_residual.png"
    q_png = FIG / "fig_quantile.png"
    if r_png.exists():
        s.shapes.add_picture(str(r_png), Inches(0.5), Inches(2.1), width=Inches(6.3))
    if q_png.exists():
        s.shapes.add_picture(str(q_png), Inches(6.85), Inches(2.15), width=Inches(5.9))
    box = s.shapes.add_textbox(Inches(0.6), Inches(6.45), Inches(12.1), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Geliştirme döneminde yüksek fiyatlarda eksik tahmin görüldü. Kilitli testte sapma tersine döndü. Yüksek fiyatlı saatlerdeki hata iki dönemde de çözülemedi.",
        size=14,
        color=NAVY_DARK,
    )
    footer(s, 7)
    notes(
        s,
        "Artık, tahminden gerçeğin çıkarılmasıyla hesaplanıyor. Pozitif değer aşırı tahmini, negatif değer eksik tahmini gösteriyor. "
        "Hata rastgele dağılmıyor. Test döneminde model ortalama olarak biraz yüksek tahmin ediyor "
        "ve yüksek fiyatlı saatlerde sapma devam ediyor. Bu yüzden model üretime hazır sayılmıyor. "
        "Geliştirme döneminde yüksek fiyatlarda eksik tahmin görüldü. Kilitli testte sapma tersine döndü ve model biraz yüksek tahmin etmeye başladı. "
        "Yüksek fiyatlı saatlerdeki hata iki dönemde de çözülemedi.",
    )

    # 8 SHAP
    s = prs.slides.add_slide(blank)
    kicker(s, "06  ·  Açıklanabilirlik")
    title(s, "Model hangi değişkenlere bakıyor?")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "SHAP ile modelin hangi değişkenlere daha çok yaslandığına bakılıyor; bu neden-sonuç kanıtlamaz. "
        "En güçlü sinyaller geçmiş fiyatlar ve gün öncesi resmi tahminler olarak görülüyor. "
        "Ay ve yılın günü gibi takvim değişkenleri tek başına fiyat sürücüsü gibi okunmuyor, çünkü birbirine bağlılar.",
        size=16,
        color=NAVY_DARK,
    )
    s.shapes.add_picture(str(shap_png), Inches(1.1), Inches(2.05), width=Inches(11.0))
    footer(s, 8)
    notes(
        s,
        "SHAP ile modelin hangi değişkenlere daha çok yaslandığına bakılıyor; bu neden-sonuç kanıtlamaz. "
        "En güçlü sinyaller geçmiş fiyatlar ve gün öncesi resmi tahminler olarak görülüyor. "
        "Ay ve yılın günü gibi takvim değişkenleri tek başına fiyat sürücüsü gibi okunmuyor, çünkü birbirine bağlılar.",
    )

    # 9 Leakage
    s = prs.slides.add_slide(blank)
    kicker(s, "07  ·  Sızıntı ve doğrulama")
    title(s, "Bu sonuca güvenebilir miyiz?")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(12.1), Inches(0.85))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Bu sonuca güvenmek için bilgi sızıntısı kapatıldı. Veri zamana göre bölündü ve satırlar karıştırılmadı. "
        "Aynı saatin gerçekleşen değerleri ve gerçekleşen piyasa fiyatı özellik olarak kullanılmadı. "
        "Model walk-forward ile seçildi; test seti ayar için kullanılmadı ve yalnızca bir kez ölçüldü.",
        size=16,
        color=NAVY_DARK,
    )
    split = FIG / "fig_split.png"
    if split.exists():
        s.shapes.add_picture(str(split), Inches(0.55), Inches(2.05), width=Inches(7.4))
    checks = [
        "Sızıntı denetimi geçti",
        "Kronolojik bölme",
        "Karıştırma yok",
        "Gerçekleşen piyasa fiyatı yok",
        "Aynı saatteki gerçekleşmeler yok",
        "Gelecek bilgi yok",
        "Test ayar için kullanılmadı",
        "Walk-forward doğrulama",
    ]
    for i, item in enumerate(checks):
        y = Inches(2.1 + i * 0.52)
        sh = rect(s, Inches(8.15), y, Inches(4.55), Inches(0.46), GREEN_BG)
        add_text(sh.text_frame, "  ✓  " + item, size=13, bold=True, color=GREEN)
    footer(s, 9)
    notes(
        s,
        "Bu sonuca güvenmek için bilgi sızıntısı kapatıldı. Veri zamana göre bölündü ve satırlar karıştırılmadı. "
        "Aynı saatin gerçekleşen değerleri ve gerçekleşen piyasa fiyatı özellik olarak kullanılmadı. "
        "Model walk-forward ile seçildi; test seti ayar için kullanılmadı ve yalnızca bir kez ölçüldü.",
    )

    # 10 Close
    s = prs.slides.add_slide(blank)
    kicker(s, "08  ·  Sonuç ve sınırlılıklar")
    title(s, "Sonuç")
    box = s.shapes.add_textbox(Inches(0.6), Inches(1.15), Inches(12.1), Inches(1.35))
    tf = box.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Sızıntısız bir tahmin hattı kuruldu ve dünün aynı saatini kopyalayan Naive taban geçildi. "
        "Model olarak Ridge, METHOD_B ve AR(1) birleşimi seçildi. "
        "Kilitli testte ortalama mutlak hata yaklaşık yüzde 34 düşürüldü. "
        "Yüksek fiyatlı saatlerdeki hata çözülemedi; veri 2014–2018 aralığında ve test tek bir kilitli dönemde ölçüldü.",
        size=17,
        color=NAVY_DARK,
    )
    goods = [
        ("Sızıntısız hat", "Kronolojik bölme, walk-forward"),
        ("Seçilen model", "Ridge + METHOD_B + AR(1)"),
        ("Kilitli MAE", "3,99 €/MWh  ·  Naive’e göre ≈ %34"),
    ]
    for i, (h, b) in enumerate(goods):
        card(s, Inches(0.6 + i * 4.1), Inches(2.7), Inches(3.9), Inches(1.25), h, b, hl=(i == 2))
    w1 = rect(s, Inches(0.6), Inches(4.2), Inches(6.0), Inches(1.15), WARN_BG)
    tf = w1.text_frame
    tf.word_wrap = True
    add_text(tf, "Sınırlılık", size=13, bold=True, color=WARN)
    add_para(tf, "Yüksek fiyatlı saatlerde hata devam ediyor.", size=15, color=NAVY_DARK, space_before=4)
    w2 = rect(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(1.15), WARN_BG)
    tf = w2.text_frame
    tf.word_wrap = True
    add_text(tf, "Sınırlılık", size=13, bold=True, color=WARN)
    add_para(tf, "Veri 2014–2018 · test tek kilitli dönem.", size=15, color=NAVY_DARK, space_before=4)
    fin = rect(s, Inches(0.6), Inches(5.55), Inches(12.1), Inches(1.15), GREEN_BG)
    tf = fin.text_frame
    tf.word_wrap = True
    add_text(
        tf,
        "Sızıntısız zaman serisi modeli kuruldu ve dünün aynı saatini kopyalayan tabandan belirgin şekilde daha iyi gün öncesi elektrik fiyat tahmini alındı.",
        size=17,
        bold=True,
        color=GREEN,
    )
    footer(s, 10)
    notes(
        s,
        "Sızıntısız bir tahmin hattı kuruldu ve dünün aynı saatini kopyalayan Naive taban geçildi. "
        "Model olarak Ridge, METHOD_B ve AR(1) birleşimi seçildi. "
        "Kilitli testte ortalama mutlak hata yaklaşık yüzde 34 düşürüldü. "
        "Yüksek fiyatlı saatlerdeki hata çözülemedi; veri 2014–2018 aralığında ve test tek bir kilitli dönemde ölçüldü.",
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    path = build()
    print(path)
